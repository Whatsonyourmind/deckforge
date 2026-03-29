"""PptxRenderer -- orchestrates rendering of a complete PPTX from IR + layout + theme."""

from __future__ import annotations

import io
import logging
from typing import TYPE_CHECKING

from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from deckforge.rendering.element_renderers import render_element
from deckforge.rendering.slide_renderers import render_finance_slide
from deckforge.rendering.utils import (
    get_blank_layout,
    set_slide_background,
    set_transition,
)

if TYPE_CHECKING:
    from deckforge.ir.presentation import Presentation
    from deckforge.layout.types import LayoutResult
    from deckforge.themes.types import ResolvedTheme

logger = logging.getLogger(__name__)


class PptxRenderer:
    """Renders a complete PPTX file from IR presentation + layout results + theme.

    The renderer creates one PPTX slide per LayoutResult, applies backgrounds,
    transitions, and speaker notes, then dispatches element rendering to the
    ELEMENT_RENDERERS registry.
    """

    def render(
        self,
        presentation: Presentation,
        layout_results: list[LayoutResult],
        theme: ResolvedTheme,
    ) -> bytes:
        """Render a complete PPTX from IR + layout results + theme.

        Args:
            presentation: The IR Presentation model.
            layout_results: List of LayoutResults from the layout engine.
            theme: Resolved theme for styling.

        Returns:
            Raw bytes of the generated .pptx file.
        """
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = get_blank_layout(prs)

        for layout_result in layout_results:
            slide = prs.slides.add_slide(blank_layout)
            ir_slide = layout_result.slide

            # 1. Set background color
            self._apply_background(slide, ir_slide, theme)

            # 2. Render elements
            self._render_elements(slide, ir_slide, theme)

            # 3. Apply transition
            self._apply_transition(slide, ir_slide)

            # 4. Apply speaker notes
            self._apply_speaker_notes(slide, ir_slide)

        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    def _apply_background(self, slide, ir_slide, theme: ResolvedTheme) -> None:
        """Set slide background from theme slide master or default theme background."""
        slide_type = ir_slide.slide_type
        if hasattr(slide_type, "value"):
            slide_type = slide_type.value

        # Look up slide master for this slide type
        master = theme.slide_masters.get(slide_type)
        if master:
            bg_color = master.background
        else:
            bg_color = theme.colors.background

        set_slide_background(slide, bg_color)

    def _render_elements(self, slide, ir_slide, theme: ResolvedTheme) -> None:
        """Render all elements on the slide.

        Finance slide types are dispatched to FINANCE_SLIDE_RENDERERS for
        full-slide rendering. Non-finance slides use the element-by-element path.
        """
        # Check if this is a finance slide type -- if so, the finance renderer
        # handles the entire slide (title, tables, charts, positioning).
        if render_finance_slide(slide, ir_slide, theme):
            return

        for element in ir_slide.elements:
            position = element.position
            if position is None:
                logger.debug(
                    "Skipping element %s -- no position assigned",
                    getattr(element, "type", "unknown"),
                )
                continue

            # Skip background elements (handled at slide level)
            element_type = element.type
            if hasattr(element_type, "value"):
                element_type = element_type.value
            if element_type == "background":
                continue

            try:
                render_element(slide, element, position, theme)
            except Exception:
                logger.exception(
                    "Failed to render element type=%s",
                    getattr(element, "type", "unknown"),
                )

    def _apply_transition(self, slide, ir_slide) -> None:
        """Apply slide transition if specified in the IR."""
        from deckforge.ir.enums import Transition

        transition = ir_slide.transition
        if transition is not None and transition != Transition.NONE:
            set_transition(slide, transition)

    def _apply_speaker_notes(self, slide, ir_slide) -> None:
        """Add speaker notes to the slide if present in the IR."""
        notes = ir_slide.speaker_notes
        if notes:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes


__all__ = ["PptxRenderer"]
