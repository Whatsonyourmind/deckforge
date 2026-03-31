#!/usr/bin/env python3
"""Generate 5 demo decks that specifically exercise chart and table rendering.

Each deck focuses on a different visualization type:
1. bar_chart_demo    - Bar chart slide
2. line_chart_demo   - Line chart slide
3. pie_chart_demo    - Pie chart slide
4. table_demo        - Table slide
5. mixed_demo        - Mixed slide (text + chart)

The IRs use the simplified format and are normalized via normalize_ir()
before rendering, validating the full normalization -> render pipeline.

Usage:
    cd SlideMaker
    pip install -e .
    python scripts/generate_chart_demos.py
"""

from __future__ import annotations

import json
import logging
import sys
import time
import traceback
from pathlib import Path

# ---------------------------------------------------------------------------
# Setup
# ---------------------------------------------------------------------------
PROJECT_ROOT = Path(__file__).resolve().parent.parent
SRC_DIR = PROJECT_ROOT / "src"
DEMOS_DIR = PROJECT_ROOT / "demos"
OUTPUT_DIR = DEMOS_DIR / "output"

if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

logging.basicConfig(level=logging.WARNING, format="%(levelname)s: %(name)s: %(message)s")
logger = logging.getLogger("generate_chart_demos")
logger.setLevel(logging.INFO)


# ---------------------------------------------------------------------------
# Demo IR definitions (simplified format -- tests normalization)
# ---------------------------------------------------------------------------

CHART_DEMOS: dict[str, dict] = {
    "bar_chart_demo": {
        "schema_version": "1.0",
        "metadata": {
            "title": "Quarterly Revenue by Region",
            "author": "DeckForge Chart Tests",
        },
        "theme": "corporate-blue",
        "slides": [
            {
                "slide_type": "title",
                "elements": [
                    {"type": "text", "content": "Quarterly Revenue by Region", "role": "title"},
                    {"type": "text", "content": "Q4 2025 Performance Analysis", "role": "subtitle"},
                ],
            },
            {
                "slide_type": "chart",
                "elements": [
                    {"type": "text", "content": "Revenue by Region ($M)", "role": "title"},
                    {
                        "type": "chart",
                        "chart_type": "bar",
                        "data": {
                            "categories": ["North America", "Europe", "Asia-Pacific", "Latin America"],
                            "series": [
                                {"name": "Q3 2025", "values": [42.5, 28.3, 18.7, 8.2]},
                                {"name": "Q4 2025", "values": [48.1, 31.6, 22.4, 9.8]},
                            ],
                        },
                    },
                    {
                        "type": "text",
                        "content": "All regions showed positive growth. North America led with 13% QoQ growth.",
                        "role": "body",
                    },
                ],
                "speaker_notes": "North America remains our strongest region with 43% of total revenue.",
            },
            {
                "slide_type": "closing",
                "elements": [
                    {"type": "text", "content": "Strong Regional Performance", "role": "title"},
                    {
                        "type": "text",
                        "content": "All four regions exceeded targets in Q4 2025.",
                        "role": "body",
                    },
                ],
            },
        ],
    },
    "line_chart_demo": {
        "schema_version": "1.0",
        "metadata": {
            "title": "User Growth Trajectory",
            "author": "DeckForge Chart Tests",
        },
        "theme": "modern-gradient",
        "slides": [
            {
                "slide_type": "title",
                "elements": [
                    {"type": "text", "content": "User Growth Trajectory", "role": "title"},
                    {"type": "text", "content": "Monthly Active Users 2024-2026", "role": "subtitle"},
                ],
            },
            {
                "slide_type": "chart",
                "elements": [
                    {"type": "text", "content": "MAU Growth (Thousands)", "role": "title"},
                    {
                        "type": "chart",
                        "chart_type": "line",
                        "data": {
                            "categories": [
                                "Jan'25", "Feb'25", "Mar'25", "Apr'25", "May'25", "Jun'25",
                                "Jul'25", "Aug'25", "Sep'25", "Oct'25", "Nov'25", "Dec'25",
                            ],
                            "series": [
                                {
                                    "name": "MAU (K)",
                                    "values": [12, 15, 19, 24, 31, 38, 47, 58, 72, 89, 108, 134],
                                },
                            ],
                        },
                    },
                ],
                "speaker_notes": "11x growth in 12 months, driven by viral referral loop.",
            },
            {
                "slide_type": "closing",
                "elements": [
                    {"type": "text", "content": "Explosive Growth Continues", "role": "title"},
                    {
                        "type": "text",
                        "content": "On track to reach 500K MAU by Q4 2026.",
                        "role": "body",
                    },
                ],
            },
        ],
    },
    "pie_chart_demo": {
        "schema_version": "1.0",
        "metadata": {
            "title": "Revenue Mix Analysis",
            "author": "DeckForge Chart Tests",
        },
        "theme": "executive-dark",
        "slides": [
            {
                "slide_type": "title",
                "elements": [
                    {"type": "text", "content": "Revenue Mix Analysis", "role": "title"},
                    {"type": "text", "content": "Product Line Contribution FY2025", "role": "subtitle"},
                ],
            },
            {
                "slide_type": "chart",
                "elements": [
                    {"type": "text", "content": "Revenue by Product Line", "role": "title"},
                    {
                        "type": "chart",
                        "chart_type": "pie",
                        "data": {
                            "categories": ["Enterprise SaaS", "SMB SaaS", "Professional Services", "Marketplace", "Other"],
                            "series": [
                                {"name": "Revenue", "values": [42, 28, 15, 10, 5]},
                            ],
                        },
                    },
                    {
                        "type": "text",
                        "content": "Enterprise SaaS represents 42% of revenue, up from 35% last year.",
                        "role": "body",
                    },
                ],
                "speaker_notes": "The shift toward enterprise is intentional -- higher ACV, lower churn, better margins.",
            },
            {
                "slide_type": "closing",
                "elements": [
                    {"type": "text", "content": "Enterprise-Led Growth", "role": "title"},
                    {
                        "type": "text",
                        "content": "Enterprise SaaS is the growth engine. Target: 55% by FY2027.",
                        "role": "body",
                    },
                ],
            },
        ],
    },
    "table_demo": {
        "schema_version": "1.0",
        "metadata": {
            "title": "Competitive Landscape Analysis",
            "author": "DeckForge Chart Tests",
        },
        "theme": "finance-pro",
        "slides": [
            {
                "slide_type": "title",
                "elements": [
                    {"type": "text", "content": "Competitive Landscape", "role": "title"},
                    {"type": "text", "content": "Market Position & Feature Comparison", "role": "subtitle"},
                ],
            },
            {
                "slide_type": "table",
                "elements": [
                    {"type": "text", "content": "Feature Comparison Matrix", "role": "title"},
                    {
                        "type": "table",
                        "data": {
                            "headers": ["Feature", "Us", "Competitor A", "Competitor B", "Competitor C"],
                            "rows": [
                                ["AI-Powered Analytics", "Yes", "Partial", "No", "Yes"],
                                ["Real-Time Dashboard", "Yes", "Yes", "Yes", "No"],
                                ["API Integration", "REST + GraphQL", "REST only", "REST only", "None"],
                                ["Pricing ($/mo)", "$99", "$149", "$199", "$79"],
                                ["Uptime SLA", "99.99%", "99.9%", "99.5%", "99.0%"],
                                ["Deployment Options", "Cloud + On-Prem", "Cloud only", "Cloud only", "Cloud only"],
                            ],
                        },
                    },
                    {
                        "type": "text",
                        "content": "We lead on 5 of 6 dimensions while being priced 34% below the category average.",
                        "role": "body",
                    },
                ],
                "speaker_notes": "The comparison is based on published specs and pricing as of Q1 2026.",
            },
            {
                "slide_type": "closing",
                "elements": [
                    {"type": "text", "content": "Best Value in Category", "role": "title"},
                    {
                        "type": "text",
                        "content": "Superior features at competitive pricing.",
                        "role": "body",
                    },
                ],
            },
        ],
    },
    "mixed_demo": {
        "schema_version": "1.0",
        "metadata": {
            "title": "Product Performance Dashboard",
            "author": "DeckForge Chart Tests",
        },
        "theme": "arctic-clean",
        "slides": [
            {
                "slide_type": "title",
                "elements": [
                    {"type": "text", "content": "Product Performance Dashboard", "role": "title"},
                    {"type": "text", "content": "Q4 2025 KPIs & Trends", "role": "subtitle"},
                ],
            },
            {
                "slide_type": "bullets",
                "elements": [
                    {"type": "text", "content": "Key Highlights", "role": "title"},
                    {
                        "type": "list",
                        "items": [
                            "Monthly active users grew 34% to 2.1M",
                            "Average session duration increased 18% to 12.4 minutes",
                            "Net Promoter Score improved from 62 to 71",
                            "Customer acquisition cost decreased 22% to $28",
                        ],
                    },
                ],
                "speaker_notes": "All four KPIs moved in the right direction this quarter.",
            },
            {
                "slide_type": "chart",
                "elements": [
                    {"type": "text", "content": "Monthly Revenue Trend", "role": "title"},
                    {
                        "type": "chart",
                        "chart_type": "bar",
                        "data": {
                            "categories": ["Oct", "Nov", "Dec"],
                            "series": [
                                {"name": "Subscriptions ($K)", "values": [320, 345, 380]},
                                {"name": "Add-ons ($K)", "values": [45, 52, 68]},
                            ],
                        },
                    },
                ],
            },
            {
                "slide_type": "table",
                "elements": [
                    {"type": "text", "content": "Product Usage Metrics", "role": "title"},
                    {
                        "type": "table",
                        "data": {
                            "headers": ["Metric", "Oct", "Nov", "Dec", "QoQ Change"],
                            "rows": [
                                ["MAU (M)", "1.57", "1.82", "2.10", "+34%"],
                                ["DAU/MAU", "0.42", "0.44", "0.47", "+12%"],
                                ["Avg Session (min)", "10.5", "11.2", "12.4", "+18%"],
                                ["Features Used/Session", "3.2", "3.5", "3.8", "+19%"],
                            ],
                        },
                    },
                ],
            },
            {
                "slide_type": "closing",
                "elements": [
                    {"type": "text", "content": "Momentum Across All Metrics", "role": "title"},
                    {
                        "type": "text",
                        "content": "Q4 demonstrated strong product-market fit acceleration.",
                        "role": "body",
                    },
                ],
            },
        ],
    },
}


# ---------------------------------------------------------------------------
# Rendering
# ---------------------------------------------------------------------------


def inspect_pptx(filepath: Path) -> dict:
    """Inspect a PPTX and return metrics."""
    from pptx import Presentation as PptxPresentation
    from pptx.util import Emu

    info: dict = {
        "file_size_kb": 0,
        "slide_count": 0,
        "total_shapes": 0,
        "has_text": False,
        "has_charts": False,
        "has_tables": False,
        "valid": False,
        "error": None,
        "slides": [],
    }

    try:
        info["file_size_kb"] = round(filepath.stat().st_size / 1024, 1)
        prs = PptxPresentation(str(filepath))
        info["slide_count"] = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            slide_info = {
                "index": i + 1,
                "shape_count": len(slide.shapes),
                "has_table": False,
                "has_chart": False,
                "text_preview": "",
            }

            for shape in slide.shapes:
                info["total_shapes"] += 1
                if shape.has_text_frame:
                    info["has_text"] = True
                    if not slide_info["text_preview"]:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_info["text_preview"] = text[:60]
                                break
                if shape.has_table:
                    info["has_tables"] = True
                    slide_info["has_table"] = True
                if shape.has_chart:
                    info["has_charts"] = True
                    slide_info["has_chart"] = True

            info["slides"].append(slide_info)

        info["valid"] = True
    except Exception as e:
        info["error"] = str(e)

    return info


def generate_chart_demo(name: str, ir_data: dict) -> dict:
    """Generate a single chart demo PPTX from simplified IR."""
    from deckforge.ir import Presentation
    from deckforge.ir.normalize import normalize_ir
    from deckforge.workers.tasks import render_pipeline

    result = {
        "name": name,
        "status": "unknown",
        "ir_slides": len(ir_data.get("slides", [])),
        "error": None,
        "inspection": None,
        "elapsed_seconds": 0.0,
    }

    start = time.time()
    output_path = OUTPUT_DIR / f"{name}.pptx"

    try:
        # 1. Normalize simplified -> strict
        normalized = normalize_ir(ir_data)

        # 2. Validate
        presentation = Presentation.model_validate(normalized)

        # 3. Render
        pptx_bytes, qa_report = render_pipeline(presentation)

        # 4. Save
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        with open(output_path, "wb") as f:
            f.write(pptx_bytes)

        # 5. Inspect
        inspection = inspect_pptx(output_path)
        result["inspection"] = inspection
        result["qa_score"] = qa_report.score
        result["qa_grade"] = qa_report.grade
        result["status"] = "PASS" if inspection["valid"] else "FAIL"

    except Exception as e:
        result["status"] = "FAIL"
        result["error"] = f"{type(e).__name__}: {e}"
        result["traceback"] = traceback.format_exc()

    result["elapsed_seconds"] = round(time.time() - start, 2)
    return result


def main() -> None:
    print("=" * 70)
    print("DeckForge Chart Demo Generation")
    print("=" * 70)
    print(f"  Output dir: {OUTPUT_DIR}")
    print(f"  Demo count: {len(CHART_DEMOS)}")
    print()

    # Also save the simplified IRs for reference
    for name, ir_data in CHART_DEMOS.items():
        demo_dir = DEMOS_DIR / name
        demo_dir.mkdir(parents=True, exist_ok=True)
        with open(demo_dir / "ir.json", "w") as f:
            json.dump(ir_data, f, indent=2)
        print(f"  Saved IR: demos/{name}/ir.json")

    print()

    results: list[dict] = []
    total_start = time.time()

    for name, ir_data in CHART_DEMOS.items():
        print(f"--- Generating: {name} ---")
        result = generate_chart_demo(name, ir_data)
        results.append(result)

        if result["status"] == "PASS":
            insp = result["inspection"]
            print(f"  [PASS] {insp['slide_count']} slides, {insp['file_size_kb']} KB")
            print(f"         QA: {result.get('qa_score')}/{result.get('qa_grade')}")
            print(f"         Text: {insp['has_text']}, Charts: {insp['has_charts']}, Tables: {insp['has_tables']}")
            for si in insp.get("slides", []):
                flags = []
                if si.get("has_chart"):
                    flags.append("CHART")
                if si.get("has_table"):
                    flags.append("TABLE")
                flag_str = f" [{', '.join(flags)}]" if flags else ""
                print(f"         Slide {si['index']}: {si['shape_count']} shapes{flag_str}  \"{si['text_preview']}\"")
        else:
            print(f"  [FAIL] {result.get('error', 'unknown')}")
            if result.get("traceback"):
                for line in result["traceback"].strip().split("\n")[-5:]:
                    print(f"    {line}")
        print(f"         Time: {result['elapsed_seconds']}s")
        print()

    total_elapsed = round(time.time() - total_start, 2)

    passed = sum(1 for r in results if r["status"] == "PASS")
    failed = sum(1 for r in results if r["status"] == "FAIL")

    print("=" * 70)
    print("SUMMARY")
    print("=" * 70)
    print(f"  Total: {len(results)}, Passed: {passed}, Failed: {failed}")
    print(f"  Time:  {total_elapsed}s")

    if passed > 0:
        print()
        print("  Generated PPTX files:")
        for r in results:
            if r["status"] == "PASS":
                insp = r["inspection"]
                print(f"    demos/output/{r['name']}.pptx ({insp['file_size_kb']} KB, {insp['slide_count']} slides)")

    if failed > 0:
        print(f"\n  RESULT: {failed}/{len(results)} FAILED")
        sys.exit(1)
    else:
        print(f"\n  RESULT: All {passed} chart demos generated successfully!")
        sys.exit(0)


if __name__ == "__main__":
    main()
