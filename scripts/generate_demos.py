#!/usr/bin/env python3
"""Generate PPTX files from demo IR JSONs and validate the output.

The demo IR files use a simplified/shorthand format that must be transformed
into the strict Pydantic IR schema before rendering. This script:

1. Loads each demo IR JSON from demos/*/ir.json
2. Transforms the simplified format into the proper IR schema
3. Renders via the render pipeline (IR -> layout -> PPTX + QA)
4. Saves output PPTX files to demos/output/
5. Inspects each PPTX for slide count, dimensions, content, shapes
6. Reports success/failure with detailed diagnostics

Usage:
    cd SlideMaker
    pip install -e .
    python scripts/generate_demos.py
"""

from __future__ import annotations

import json
import logging
import os
import sys
import time
import traceback
from pathlib import Path

# ---------------------------------------------------------------------------
# Setup paths
# ---------------------------------------------------------------------------
PROJECT_ROOT = Path(__file__).resolve().parent.parent
SRC_DIR = PROJECT_ROOT / "src"
DEMOS_DIR = PROJECT_ROOT / "demos"
OUTPUT_DIR = DEMOS_DIR / "output"

# Ensure deckforge is importable
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

# Configure logging
logging.basicConfig(
    level=logging.WARNING,
    format="%(levelname)s: %(name)s: %(message)s",
)
logger = logging.getLogger("generate_demos")
logger.setLevel(logging.INFO)


# ---------------------------------------------------------------------------
# Slide type mapping: simplified demo names -> strict enum values
# ---------------------------------------------------------------------------
SLIDE_TYPE_MAP: dict[str, str] = {
    # Shorthand -> SlideType enum value
    "title": "title_slide",
    "bullets": "bullet_points",
    "chart": "chart_slide",
    "table": "table_slide",
    "two_column": "two_column_text",
    "team": "team_slide",
    "closing": "thank_you",
    "quote": "quote_slide",
    "executive_summary": "key_message",
    # These already match
    "title_slide": "title_slide",
    "bullet_points": "bullet_points",
    "chart_slide": "chart_slide",
    "table_slide": "table_slide",
    "two_column_text": "two_column_text",
    "team_slide": "team_slide",
    "thank_you": "thank_you",
    "comparison": "comparison",
    "timeline": "timeline",
    "process_flow": "process_flow",
    "section_divider": "section_divider",
    "key_message": "key_message",
    "quote_slide": "quote_slide",
    "stats_callout": "stats_callout",
    "agenda": "agenda",
    "image_with_caption": "image_with_caption",
    "icon_grid": "icon_grid",
    "matrix": "matrix",
    "funnel": "funnel",
    "map_slide": "map_slide",
    "appendix": "appendix",
    "q_and_a": "q_and_a",
    "org_chart": "org_chart",
    # Finance
    "dcf_summary": "dcf_summary",
    "comp_table": "comp_table",
    "waterfall_chart": "waterfall_chart",
    "deal_overview": "deal_overview",
    "returns_analysis": "returns_analysis",
    "capital_structure": "capital_structure",
    "market_landscape": "market_landscape",
    "risk_matrix": "risk_matrix",
    "investment_thesis": "investment_thesis",
}

# Role -> element type mapping for "type": "text" elements
ROLE_TO_ELEMENT_TYPE: dict[str, str] = {
    "title": "heading",
    "subtitle": "subheading",
    "body": "body_text",
    "left_header": "heading",
    "right_header": "heading",
    "left_body": "body_text",
    "right_body": "body_text",
    "attribution": "footnote",
}


def transform_chart_data(chart_type: str, data: dict) -> dict:
    """Transform chart data from simplified demo format to strict ChartUnion schema.

    Some chart types (waterfall, funnel, pie, donut) have different schemas
    than the generic category+series format used in demos.
    """
    categories = data.get("categories", [])
    series = data.get("series", [])

    if chart_type == "waterfall":
        # WaterfallChartData: categories + values (flat list, not series)
        values = series[0]["values"] if series else []
        return {
            "chart_type": "waterfall",
            "categories": categories,
            "values": values,
        }

    elif chart_type == "funnel":
        # FunnelChartData: stages + values (flat list, not categories + series)
        values = series[0]["values"] if series else []
        return {
            "chart_type": "funnel",
            "stages": categories,
            "values": values,
        }

    elif chart_type in ("pie", "donut"):
        # PieChartData / DonutChartData: labels + values
        values = series[0]["values"] if series else []
        result = {
            "chart_type": chart_type,
            "labels": categories,
            "values": values,
        }
        if chart_type == "donut":
            result["inner_radius"] = 0.5
        return result

    elif chart_type == "radar":
        # RadarChartData: axes + series
        return {
            "chart_type": "radar",
            "axes": categories,
            "series": series,
        }

    elif chart_type == "treemap":
        # TreemapChartData: labels + values
        values = series[0]["values"] if series else []
        return {
            "chart_type": "treemap",
            "labels": categories,
            "values": values,
        }

    else:
        # Generic category-based chart (bar, line, area, stacked_bar, etc.)
        return {
            "chart_type": chart_type,
            "categories": categories,
            "series": series,
        }


def transform_element(elem: dict) -> dict | None:
    """Transform a simplified element dict into the strict IR schema format.

    Handles:
      - "type": "text" + "role" -> heading/subheading/body_text with nested content
      - "type": "list" + "items" -> bullet_list with nested content
      - "type": "table" + "data" -> table with nested content
      - "type": "chart" + "chart_type" + "data" -> chart with chart_data
      - "type": "timeline" -> bullet_list (timeline items as bullet text)
    """
    elem_type = elem.get("type", "")
    role = elem.get("role", "body")

    if elem_type == "text":
        content_text = elem.get("content", "")
        ir_type = ROLE_TO_ELEMENT_TYPE.get(role, "body_text")

        if ir_type == "heading":
            return {
                "type": "heading",
                "content": {"text": content_text, "level": "h1" if role == "title" else "h2"},
            }
        elif ir_type == "subheading":
            return {
                "type": "subheading",
                "content": {"text": content_text},
            }
        elif ir_type == "footnote":
            return {
                "type": "footnote",
                "content": {"text": content_text},
            }
        else:
            return {
                "type": "body_text",
                "content": {"text": content_text},
            }

    elif elem_type == "list":
        items = elem.get("items", [])
        return {
            "type": "bullet_list",
            "content": {"items": items},
        }

    elif elem_type == "table":
        data = elem.get("data", {})
        headers = data.get("headers", [])
        rows = data.get("rows", [])
        return {
            "type": "table",
            "content": {
                "headers": headers,
                "rows": rows,
            },
        }

    elif elem_type == "chart":
        chart_type = elem.get("chart_type", "bar")
        data = elem.get("data", {})
        # Build chart_data in the ChartUnion format, handling schema differences
        chart_data = transform_chart_data(chart_type, data)
        return {
            "type": "chart",
            "chart_data": chart_data,
        }

    elif elem_type == "timeline":
        # Timeline items are not a standard element type.
        # Convert to a bullet list with formatted text.
        items_raw = elem.get("items", [])
        bullet_items = []
        for item in items_raw:
            if isinstance(item, dict):
                date = item.get("date", "")
                title = item.get("title", "")
                desc = item.get("description", "")
                bullet_items.append(f"{date}: {title} -- {desc}")
            else:
                bullet_items.append(str(item))
        return {
            "type": "bullet_list",
            "content": {"items": bullet_items},
        }

    else:
        # Unknown element type -- skip
        logger.warning("Unknown element type %r, skipping", elem_type)
        return None


def transform_slide(slide: dict) -> dict:
    """Transform a simplified slide dict into the strict IR schema format."""
    raw_type = slide.get("slide_type", "")
    mapped_type = SLIDE_TYPE_MAP.get(raw_type, raw_type)

    transformed: dict = {
        "slide_type": mapped_type,
        "elements": [],
    }

    # Preserve optional fields
    if slide.get("speaker_notes"):
        transformed["speaker_notes"] = slide["speaker_notes"]
    if slide.get("layout_hint"):
        transformed["layout_hint"] = slide["layout_hint"]
    if slide.get("transition"):
        transformed["transition"] = slide["transition"]

    # Transform elements
    for elem in slide.get("elements", []):
        new_elem = transform_element(elem)
        if new_elem is not None:
            transformed["elements"].append(new_elem)

    return transformed


def transform_ir(ir_data: dict) -> dict:
    """Transform a full simplified IR dict into the strict schema format."""
    transformed = {
        "schema_version": ir_data.get("schema_version", "1.0"),
        "metadata": ir_data.get("metadata", {}),
        "theme": ir_data.get("theme", "executive-dark"),
        "slides": [],
    }

    if ir_data.get("brand_kit"):
        transformed["brand_kit"] = ir_data["brand_kit"]
    if ir_data.get("generation_options"):
        transformed["generation_options"] = ir_data["generation_options"]

    for slide in ir_data.get("slides", []):
        transformed["slides"].append(transform_slide(slide))

    return transformed


def inspect_pptx(filepath: Path) -> dict:
    """Inspect a generated PPTX file and return detailed metrics."""
    from pptx import Presentation as PptxPresentation
    from pptx.util import Emu

    info: dict = {
        "file_size_bytes": 0,
        "file_size_kb": 0,
        "slide_count": 0,
        "width_inches": 0.0,
        "height_inches": 0.0,
        "slides": [],
        "total_shapes": 0,
        "total_text_frames": 0,
        "has_text": False,
        "has_shapes": False,
        "has_charts": False,
        "has_tables": False,
        "valid": False,
        "error": None,
    }

    try:
        info["file_size_bytes"] = filepath.stat().st_size
        info["file_size_kb"] = round(info["file_size_bytes"] / 1024, 1)

        prs = PptxPresentation(str(filepath))
        info["width_inches"] = round(prs.slide_width / Emu(914400), 2)
        info["height_inches"] = round(prs.slide_height / Emu(914400), 2)
        info["slide_count"] = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            slide_info: dict = {
                "index": i + 1,
                "shape_count": len(slide.shapes),
                "has_text": False,
                "has_table": False,
                "has_chart": False,
                "text_preview": "",
            }

            for shape in slide.shapes:
                info["total_shapes"] += 1

                if shape.has_text_frame:
                    info["total_text_frames"] += 1
                    info["has_text"] = True
                    slide_info["has_text"] = True
                    # Grab first non-empty text as preview
                    if not slide_info["text_preview"]:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_info["text_preview"] = text[:80]
                                break

                if shape.has_table:
                    info["has_tables"] = True
                    slide_info["has_table"] = True

                if shape.has_chart:
                    info["has_charts"] = True
                    slide_info["has_chart"] = True

            info["has_shapes"] = info["total_shapes"] > 0
            info["slides"].append(slide_info)

        info["valid"] = True

    except Exception as e:
        info["error"] = str(e)

    return info


def generate_demo(demo_name: str) -> dict:
    """Generate a PPTX from a demo IR and return results."""
    result = {
        "demo": demo_name,
        "status": "unknown",
        "ir_slides": 0,
        "pptx_slides": 0,
        "file_size_kb": 0,
        "qa_score": None,
        "qa_grade": None,
        "transform_warnings": [],
        "error": None,
        "inspection": None,
        "elapsed_seconds": 0.0,
    }

    ir_path = DEMOS_DIR / demo_name / "ir.json"
    output_path = OUTPUT_DIR / f"{demo_name}.pptx"

    if not ir_path.exists():
        result["status"] = "FAIL"
        result["error"] = f"IR file not found: {ir_path}"
        return result

    start = time.time()

    try:
        # 1. Load raw IR
        with open(ir_path) as f:
            raw_ir = json.load(f)
        result["ir_slides"] = len(raw_ir.get("slides", []))

        # 2. Transform to strict schema
        transformed = transform_ir(raw_ir)

        # 3. Validate against Pydantic model
        from deckforge.ir import Presentation

        presentation = Presentation.model_validate(transformed)

        # 4. Render via render_pipeline
        from deckforge.workers.tasks import render_pipeline

        pptx_bytes, qa_report = render_pipeline(presentation)
        result["qa_score"] = qa_report.score
        result["qa_grade"] = qa_report.grade

        # 5. Save output
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        with open(output_path, "wb") as f:
            f.write(pptx_bytes)

        # 6. Inspect the PPTX
        inspection = inspect_pptx(output_path)
        result["inspection"] = inspection
        result["pptx_slides"] = inspection["slide_count"]
        result["file_size_kb"] = inspection["file_size_kb"]
        result["status"] = "PASS" if inspection["valid"] else "FAIL"

    except Exception as e:
        result["status"] = "FAIL"
        result["error"] = f"{type(e).__name__}: {e}"
        result["traceback"] = traceback.format_exc()

    result["elapsed_seconds"] = round(time.time() - start, 2)
    return result


def print_result(r: dict) -> None:
    """Print formatted result for a single demo."""
    status_icon = "PASS" if r["status"] == "PASS" else "FAIL"
    print(f"  [{status_icon}] {r['demo']}")
    print(f"    IR slides:    {r['ir_slides']}")

    if r["status"] == "PASS":
        insp = r["inspection"]
        print(f"    PPTX slides:  {r['pptx_slides']}")
        print(f"    Dimensions:   {insp['width_inches']}\" x {insp['height_inches']}\"")
        print(f"    File size:    {r['file_size_kb']} KB")
        print(f"    Total shapes: {insp['total_shapes']}")
        print(f"    Text frames:  {insp['total_text_frames']}")
        print(f"    Has text:     {insp['has_text']}")
        print(f"    Has tables:   {insp['has_tables']}")
        print(f"    Has charts:   {insp['has_charts']}")
        print(f"    QA score:     {r['qa_score']}")
        print(f"    QA grade:     {r['qa_grade']}")
        print(f"    Time:         {r['elapsed_seconds']}s")

        # Show per-slide text previews
        for slide_info in insp.get("slides", []):
            shapes = slide_info["shape_count"]
            preview = slide_info.get("text_preview", "")
            flags = []
            if slide_info.get("has_table"):
                flags.append("TABLE")
            if slide_info.get("has_chart"):
                flags.append("CHART")
            flag_str = f" [{', '.join(flags)}]" if flags else ""
            print(f"      Slide {slide_info['index']}: {shapes} shapes{flag_str}  \"{preview}\"")
    else:
        print(f"    Error:        {r.get('error', 'unknown')}")
        if r.get("traceback"):
            # Show last 5 lines of traceback
            tb_lines = r["traceback"].strip().split("\n")
            for line in tb_lines[-5:]:
                print(f"      {line}")
        print(f"    Time:         {r['elapsed_seconds']}s")


def main() -> None:
    """Generate and validate all demo decks."""
    print("=" * 70)
    print("DeckForge Demo Generation & Validation")
    print("=" * 70)
    print(f"  Demos dir:  {DEMOS_DIR}")
    print(f"  Output dir: {OUTPUT_DIR}")
    print()

    # Discover demos
    demos = sorted(
        d.name
        for d in DEMOS_DIR.iterdir()
        if d.is_dir() and (d / "ir.json").exists()
    )

    if not demos:
        print("No demo IR files found!")
        sys.exit(1)

    print(f"Found {len(demos)} demos: {', '.join(demos)}")
    print()

    # Generate each demo
    results: list[dict] = []
    total_start = time.time()

    for demo_name in demos:
        print(f"--- Generating: {demo_name} ---")
        result = generate_demo(demo_name)
        results.append(result)
        print_result(result)
        print()

    total_elapsed = round(time.time() - total_start, 2)

    # Summary
    passed = sum(1 for r in results if r["status"] == "PASS")
    failed = sum(1 for r in results if r["status"] == "FAIL")

    print("=" * 70)
    print("SUMMARY")
    print("=" * 70)
    print(f"  Total demos:  {len(results)}")
    print(f"  Passed:       {passed}")
    print(f"  Failed:       {failed}")
    print(f"  Total time:   {total_elapsed}s")
    print()

    if passed > 0:
        print("  Generated PPTX files:")
        for r in results:
            if r["status"] == "PASS":
                print(f"    demos/output/{r['demo']}.pptx ({r['file_size_kb']} KB, {r['pptx_slides']} slides)")
    print()

    if failed > 0:
        print(f"  RESULT: {failed}/{len(results)} demos FAILED")
        sys.exit(1)
    else:
        print(f"  RESULT: All {passed} demos generated successfully!")
        sys.exit(0)


if __name__ == "__main__":
    main()
