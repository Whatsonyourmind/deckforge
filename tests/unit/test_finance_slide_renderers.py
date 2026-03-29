"""Tests for finance slide renderers -- comp_table, dcf_summary, waterfall, and registry."""

from __future__ import annotations

import io
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from deckforge.ir.elements.base import Position
from deckforge.ir.elements.data import (
    ChartElement,
    TableContent,
    TableElement,
)
from deckforge.ir.elements.text import (
    BodyTextContent,
    BodyTextElement,
    HeadingContent,
    HeadingElement,
)
from deckforge.themes.types import (
    ResolvedTheme,
    ThemeColors,
    ThemeSpacing,
    ThemeTypography,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture()
def theme() -> ResolvedTheme:
    """Minimal resolved theme for testing."""
    return ResolvedTheme(
        name="test-finance",
        description="Test theme for finance renderers",
        colors=ThemeColors(
            primary="#0A1E3D",
            secondary="#1A3A5C",
            accent="#FF6B35",
            background="#FFFFFF",
            surface="#F5F5F5",
            text_primary="#1A1A1A",
            text_secondary="#4A4A4A",
            text_muted="#8A8A8A",
            positive="#28A745",
            negative="#DC3545",
            warning="#FFC107",
        ),
        typography=ThemeTypography(
            heading_family="Arial",
            body_family="Calibri",
            mono_family="Consolas",
        ),
        spacing=ThemeSpacing(
            margin_top=0.5,
            margin_bottom=0.5,
            margin_left=0.75,
            margin_right=0.75,
            gutter=0.3,
            element_gap=0.2,
            section_gap=0.5,
        ),
    )


@pytest.fixture()
def pptx_slide():
    """Create a blank PPTX slide for testing."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    return slide


def _make_comp_table_slide():
    """Create a CompTableSlide with sample financial data."""
    from deckforge.ir.slides.finance import CompTableSlide

    table_elem = TableElement(
        content=TableContent(
            headers=["Company", "EV/EBITDA", "P/E", "Market Cap", "Revenue Growth"],
            rows=[
                ["Company A", 12.5, 18.3, 4200000000, 0.15],
                ["Company B", 10.2, 14.7, 2800000000, 0.22],
                ["Company C", 15.8, 22.1, 6100000000, 0.08],
                ["Company D", 11.0, 16.5, 3500000000, 0.18],
            ],
        ),
        position=Position(x=0.75, y=1.5, width=11.8, height=4.5),
    )
    heading = HeadingElement(
        content=HeadingContent(text="Comparable Companies Analysis", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    return CompTableSlide(elements=[heading, table_elem])


def _make_dcf_slide():
    """Create a DcfSummarySlide with assumptions and sensitivity data."""
    from deckforge.ir.slides.finance import DcfSummarySlide

    # Assumptions as a table element
    assumptions_table = TableElement(
        content=TableContent(
            headers=["Assumption", "Value"],
            rows=[
                ["WACC", 0.10],
                ["Terminal Growth Rate", 0.025],
                ["Projection Period", 5],
                ["Terminal Value Method", "Gordon Growth"],
            ],
        ),
        position=Position(x=0.75, y=1.5, width=5.0, height=3.0),
    )

    # Sensitivity matrix as a table element with numeric-looking headers
    sensitivity_table = TableElement(
        content=TableContent(
            headers=["Discount Rate", "1.5%", "2.0%", "2.5%", "3.0%"],
            rows=[
                ["8.0%", 145, 152, 160, 170],
                ["9.0%", 130, 136, 142, 150],
                ["10.0%", 118, 123, 128, 135],
                ["11.0%", 108, 112, 116, 122],
            ],
        ),
        position=Position(x=6.5, y=1.5, width=6.0, height=3.0),
    )

    heading = HeadingElement(
        content=HeadingContent(text="DCF Valuation Summary", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )

    return DcfSummarySlide(
        elements=[heading, assumptions_table, sensitivity_table],
        discount_rate_range=[0.08, 0.09, 0.10, 0.11],
        terminal_growth_range=[0.015, 0.020, 0.025, 0.030],
    )


def _make_waterfall_slide():
    """Create a WaterfallChartSlide with chart data."""
    from deckforge.ir.charts.types import WaterfallChartData
    from deckforge.ir.slides.finance import WaterfallChartSlide

    chart_elem = ChartElement(
        chart_data=WaterfallChartData(
            categories=["Revenue", "COGS", "Gross Profit", "OpEx", "EBITDA", "Net Income"],
            values=[100, -40, 60, -25, 35, 28],
            title="Bridge to Net Income",
        ),
        position=Position(x=0.75, y=1.5, width=11.8, height=5.0),
    )
    heading = HeadingElement(
        content=HeadingContent(text="Earnings Bridge", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    return WaterfallChartSlide(elements=[heading, chart_elem], show_running_total=True)


# ── CompTableRenderer Tests ──────────────────────────────────────────────────


class TestCompTableRenderer:
    def test_produces_table_shape(self, pptx_slide, theme):
        from deckforge.rendering.slide_renderers.comp_table import CompTableRenderer

        renderer = CompTableRenderer()
        ir_slide = _make_comp_table_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Should have at least one table shape
        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) >= 1, "CompTableRenderer should produce at least one table shape"

    def test_table_has_formatted_numbers(self, pptx_slide, theme):
        from deckforge.rendering.slide_renderers.comp_table import CompTableRenderer

        renderer = CompTableRenderer()
        ir_slide = _make_comp_table_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Find the main data table
        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert table_shapes, "Should have a table"
        table = table_shapes[0].table

        # Collect all cell text
        all_text = []
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                all_text.append(table.cell(row_idx, col_idx).text)

        joined = " ".join(all_text)
        # Should have formatted multiples (e.g., "12.5x")
        assert "x" in joined, "Should contain multiple formatted values (e.g., 12.5x)"
        # Should have formatted currency (e.g., "$4.2B")
        assert "$" in joined, "Should contain currency formatted values"

    def test_median_row_highlighted(self, pptx_slide, theme):
        from deckforge.rendering.slide_renderers.comp_table import CompTableRenderer

        renderer = CompTableRenderer()
        ir_slide = _make_comp_table_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert table_shapes
        table = table_shapes[0].table

        # The last row should be the median footer row
        last_row_idx = len(table.rows) - 1
        cell_text = table.cell(last_row_idx, 0).text
        assert "median" in cell_text.lower() or "mean" in cell_text.lower(), \
            f"Last row should be median/mean, got: {cell_text}"

    def test_numeric_columns_right_aligned(self, pptx_slide, theme):
        from pptx.enum.text import PP_ALIGN

        from deckforge.rendering.slide_renderers.comp_table import CompTableRenderer

        renderer = CompTableRenderer()
        ir_slide = _make_comp_table_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert table_shapes
        table = table_shapes[0].table

        # Check a numeric cell (row 1, col 1 = EV/EBITDA value)
        cell = table.cell(1, 1)
        for para in cell.text_frame.paragraphs:
            if para.alignment is not None:
                assert para.alignment == PP_ALIGN.RIGHT, "Numeric columns should be right-aligned"


# ── DcfSummaryRenderer Tests ─────────────────────────────────────────────────


class TestDcfSummaryRenderer:
    def test_produces_at_least_two_tables(self, pptx_slide, theme):
        from deckforge.rendering.slide_renderers.dcf_summary import DcfSummaryRenderer

        renderer = DcfSummaryRenderer()
        ir_slide = _make_dcf_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) >= 2, \
            f"DcfSummaryRenderer should produce at least 2 tables (assumptions + sensitivity), got {len(table_shapes)}"

    def test_sensitivity_matrix_has_gradient_colors(self, pptx_slide, theme):
        from deckforge.rendering.slide_renderers.dcf_summary import DcfSummaryRenderer

        renderer = DcfSummaryRenderer()
        ir_slide = _make_dcf_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Find the sensitivity table (the one with more columns)
        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) >= 2

        # Get the larger table (sensitivity matrix)
        sensitivity_table = max(table_shapes, key=lambda s: len(s.table.columns))
        table = sensitivity_table.table

        # Check that data cells have fill colors applied (not all the same)
        fills = set()
        for row_idx in range(1, len(table.rows)):  # Skip header
            for col_idx in range(1, len(table.columns)):  # Skip row header
                cell = table.cell(row_idx, col_idx)
                if cell.fill.type is not None:
                    fills.add(str(cell.fill.fore_color.rgb))

        assert len(fills) >= 2, "Sensitivity matrix should have varied gradient colors"


# ── WaterfallSlideRenderer Tests ─────────────────────────────────────────────


class TestWaterfallSlideRenderer:
    @patch("deckforge.rendering.chart_renderers.waterfall.WaterfallChartRenderer.render")
    def test_produces_slide_with_chart(self, mock_render, pptx_slide, theme):
        """WaterfallSlideRenderer should delegate chart rendering to WaterfallChartRenderer."""
        from deckforge.rendering.slide_renderers.waterfall_slide import WaterfallSlideRenderer

        renderer = WaterfallSlideRenderer()
        ir_slide = _make_waterfall_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # The renderer should have called WaterfallChartRenderer.render
        assert mock_render.called, "Should delegate to WaterfallChartRenderer"

    def test_has_title(self, pptx_slide, theme):
        """WaterfallSlideRenderer should add a title text box."""
        from deckforge.rendering.slide_renderers.waterfall_slide import WaterfallSlideRenderer

        # Mock the chart renderer to avoid Plotly image generation in tests
        with patch("deckforge.rendering.chart_renderers.waterfall.WaterfallChartRenderer.render"):
            renderer = WaterfallSlideRenderer()
            ir_slide = _make_waterfall_slide()
            renderer.render(pptx_slide, ir_slide, theme)

        # Should have a text shape with the title
        text_shapes = [s for s in pptx_slide.shapes if s.has_text_frame]
        title_texts = [s.text_frame.text for s in text_shapes]
        assert any("Earnings Bridge" in t or "Bridge" in t for t in title_texts), \
            f"Should have title text box, found: {title_texts}"


# ── Registry Tests ───────────────────────────────────────────────────────────


class TestFinanceSlideRegistry:
    def test_registry_has_comp_table(self):
        from deckforge.rendering.slide_renderers import FINANCE_SLIDE_RENDERERS
        from deckforge.rendering.slide_renderers.comp_table import CompTableRenderer

        assert "comp_table" in FINANCE_SLIDE_RENDERERS
        assert isinstance(FINANCE_SLIDE_RENDERERS["comp_table"], CompTableRenderer)

    def test_registry_has_all_9_types(self):
        from deckforge.rendering.slide_renderers import FINANCE_SLIDE_RENDERERS

        expected_types = {
            "comp_table",
            "dcf_summary",
            "waterfall_chart",
            "deal_overview",
            "returns_analysis",
            "capital_structure",
            "market_landscape",
            "investment_thesis",
            "risk_matrix",
        }
        assert set(FINANCE_SLIDE_RENDERERS.keys()) == expected_types, \
            f"Missing types: {expected_types - set(FINANCE_SLIDE_RENDERERS.keys())}"

    def test_render_finance_slide_dispatches(self, pptx_slide, theme):
        from deckforge.rendering.slide_renderers import render_finance_slide

        ir_slide = _make_comp_table_slide()
        result = render_finance_slide(pptx_slide, ir_slide, theme)
        assert result is True, "render_finance_slide should return True for comp_table"

    def test_render_finance_slide_returns_false_for_unknown(self, pptx_slide, theme):
        from deckforge.rendering.slide_renderers import render_finance_slide

        # Create a mock slide with unknown type
        mock_slide = MagicMock()
        mock_slide.slide_type = "title_slide"
        result = render_finance_slide(pptx_slide, mock_slide, theme)
        assert result is False, "render_finance_slide should return False for unknown types"
