"""Tests for static chart renderers -- Plotly-based PNG image charts on PPTX slides."""

from __future__ import annotations

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from deckforge.ir.charts.types import (
    FootballFieldChartData,
    FunnelChartData,
    GanttChartData,
    GanttTask,
    HeatmapChartData,
    SankeyChartData,
    SankeyLink,
    SensitivityTableData,
    SunburstChartData,
    TreemapChartData,
    TornadoChartData,
    WaterfallChartData,
)
from deckforge.ir.elements.base import Position
from deckforge.themes.types import (
    ResolvedTheme,
    ThemeColors,
    ThemeSpacing,
    ThemeTypography,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture
def pptx_slide():
    """Create a blank PPTX presentation slide for chart testing."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    return slide


@pytest.fixture
def theme():
    """Minimal resolved theme for chart formatting."""
    return ResolvedTheme(
        name="test-theme",
        description="Test theme",
        colors=ThemeColors(
            primary="#2E86AB",
            secondary="#A23B72",
            accent="#F18F01",
            background="#1A1A2E",
            surface="#16213E",
            text_primary="#EAEAEA",
            text_secondary="#B0B0B0",
            text_muted="#808080",
            positive="#4CAF50",
            negative="#F44336",
            warning="#FF9800",
        ),
        typography=ThemeTypography(
            heading_family="Arial",
            body_family="Calibri",
            mono_family="Consolas",
        ),
        spacing=ThemeSpacing(
            margin_top=0.5,
            margin_bottom=0.5,
            margin_left=0.75,
            margin_right=0.75,
            gutter=0.3,
            element_gap=0.2,
            section_gap=0.5,
        ),
        chart_colors=["#2E86AB", "#A23B72", "#F18F01", "#4CAF50", "#FF9800", "#9C27B0"],
    )


@pytest.fixture
def position():
    """Standard chart position on a slide."""
    return Position(x=1.0, y=1.5, width=11.0, height=5.0)


# ── StaticChartRenderer base tests ────────────────────────────────────────────


class TestStaticChartRendererBase:
    def test_render_to_image_returns_png_bytes(self, theme, position):
        """StaticChartRenderer._render_to_image() returns PNG bytes (length > 0)."""
        from deckforge.rendering.chart_renderers.static_base import StaticChartRenderer

        import plotly.graph_objects as go

        # Create a concrete subclass for testing
        class TestRenderer(StaticChartRenderer):
            def _build_figure(self, chart_data, theme):
                fig = go.Figure(data=[go.Bar(x=["A", "B"], y=[1, 2])])
                return fig

        renderer = TestRenderer()
        fig = go.Figure(data=[go.Bar(x=["A", "B"], y=[1, 2])])
        img_bytes = renderer._render_to_image(fig, position)
        assert isinstance(img_bytes, bytes)
        assert len(img_bytes) > 0
        # PNG magic bytes
        assert img_bytes[:4] == b"\x89PNG"


# ── Waterfall Chart ──────────────────────────────────────────────────────────


class TestWaterfallChartRenderer:
    def test_waterfall_produces_picture_shape(self, pptx_slide, theme, position):
        """WaterfallChartRenderer.render() produces a slide with exactly 1 picture shape."""
        from deckforge.rendering.chart_renderers.waterfall import WaterfallChartRenderer

        renderer = WaterfallChartRenderer()
        chart_data = WaterfallChartData(
            categories=["Revenue", "Costs", "Profit"],
            values=[100, -30, 70],
            title="Bridge",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        # Should have exactly 1 shape and it should be a picture (not a rectangle)
        assert len(pptx_slide.shapes) == 1
        shape = pptx_slide.shapes[0]
        assert shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE


# ── Heatmap Chart ────────────────────────────────────────────────────────────


class TestHeatmapChartRenderer:
    def test_heatmap_produces_picture_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers.heatmap import HeatmapChartRenderer

        renderer = HeatmapChartRenderer()
        chart_data = HeatmapChartData(
            x_labels=["Q1", "Q2", "Q3"],
            y_labels=["Product A", "Product B"],
            data=[[10, 20, 30], [40, 50, 60]],
            title="Sales Heatmap",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1
        assert pptx_slide.shapes[0].shape_type == 13


# ── Sankey Chart ─────────────────────────────────────────────────────────────


class TestSankeyChartRenderer:
    def test_sankey_produces_picture_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers.sankey import SankeyChartRenderer

        renderer = SankeyChartRenderer()
        chart_data = SankeyChartData(
            nodes=["Revenue", "COGS", "Gross Profit", "OpEx", "Net Income"],
            links=[
                SankeyLink(source="Revenue", target="COGS", value=40),
                SankeyLink(source="Revenue", target="Gross Profit", value=60),
                SankeyLink(source="Gross Profit", target="OpEx", value=30),
                SankeyLink(source="Gross Profit", target="Net Income", value=30),
            ],
            title="P&L Flow",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1
        assert pptx_slide.shapes[0].shape_type == 13


# ── Gantt Chart ──────────────────────────────────────────────────────────────


class TestGanttChartRenderer:
    def test_gantt_produces_picture_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers.gantt import GanttChartRenderer

        renderer = GanttChartRenderer()
        chart_data = GanttChartData(
            tasks=[
                GanttTask(name="Design", start="2026-01-01", end="2026-01-15"),
                GanttTask(name="Build", start="2026-01-10", end="2026-02-01"),
                GanttTask(name="Test", start="2026-02-01", end="2026-02-15"),
            ],
            title="Project Timeline",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1
        assert pptx_slide.shapes[0].shape_type == 13


# ── Football Field Chart ─────────────────────────────────────────────────────


class TestFootballFieldChartRenderer:
    def test_football_field_produces_picture_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers.football_field import (
            FootballFieldChartRenderer,
        )

        renderer = FootballFieldChartRenderer()
        chart_data = FootballFieldChartData(
            categories=["DCF", "Comps", "Precedents"],
            low_values=[80, 90, 85],
            high_values=[120, 130, 125],
            mid_values=[100, 110, 105],
            title="Valuation Range",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1
        assert pptx_slide.shapes[0].shape_type == 13


# ── Sensitivity Chart ────────────────────────────────────────────────────────


class TestSensitivityChartRenderer:
    def test_sensitivity_produces_picture_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers.sensitivity import (
            SensitivityChartRenderer,
        )

        renderer = SensitivityChartRenderer()
        chart_data = SensitivityTableData(
            row_header="WACC",
            col_header="Growth Rate",
            row_values=[8.0, 9.0, 10.0],
            col_values=[1.0, 2.0, 3.0],
            data=[[120, 110, 100], [115, 105, 95], [110, 100, 90]],
            title="DCF Sensitivity",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1
        assert pptx_slide.shapes[0].shape_type == 13


# ── Funnel Chart ─────────────────────────────────────────────────────────────


class TestFunnelChartRenderer:
    def test_funnel_produces_picture_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers.funnel import FunnelChartRenderer

        renderer = FunnelChartRenderer()
        chart_data = FunnelChartData(
            stages=["Leads", "Qualified", "Proposal", "Won"],
            values=[1000, 500, 200, 100],
            title="Sales Funnel",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1
        assert pptx_slide.shapes[0].shape_type == 13


# ── Treemap Chart ────────────────────────────────────────────────────────────


class TestTreemapChartRenderer:
    def test_treemap_produces_picture_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers.treemap import TreemapChartRenderer

        renderer = TreemapChartRenderer()
        chart_data = TreemapChartData(
            labels=["Total", "A", "B", "C"],
            values=[0, 50, 30, 20],
            parents=["", "Total", "Total", "Total"],
            title="Revenue Breakdown",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1
        assert pptx_slide.shapes[0].shape_type == 13


# ── Tornado Chart ────────────────────────────────────────────────────────────


class TestTornadoChartRenderer:
    def test_tornado_produces_picture_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers.tornado import TornadoChartRenderer

        renderer = TornadoChartRenderer()
        chart_data = TornadoChartData(
            categories=["Revenue Growth", "WACC", "Terminal Multiple"],
            low_values=[-10, -8, -5],
            high_values=[10, 8, 5],
            base_value=0,
            title="Sensitivity Tornado",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1
        assert pptx_slide.shapes[0].shape_type == 13


# ── Sunburst Chart ───────────────────────────────────────────────────────────


class TestSunburstChartRenderer:
    def test_sunburst_produces_picture_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers.sunburst import SunburstChartRenderer

        renderer = SunburstChartRenderer()
        chart_data = SunburstChartData(
            labels=["Total", "A", "B", "A1", "A2"],
            parents=["", "Total", "Total", "A", "A"],
            values=[0, 0, 30, 40, 20],
            title="Org Breakdown",
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1
        assert pptx_slide.shapes[0].shape_type == 13


# ── Registry tests ───────────────────────────────────────────────────────────


class TestStaticChartRegistry:
    def test_waterfall_is_not_placeholder(self):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS
        from deckforge.rendering.chart_renderers.placeholder import PlaceholderChartRenderer

        assert not isinstance(CHART_RENDERERS["waterfall"], PlaceholderChartRenderer)

    def test_all_10_types_are_not_placeholder(self):
        """All 10 previously-placeholder types now return non-Placeholder renderers."""
        from deckforge.rendering.chart_renderers import CHART_RENDERERS
        from deckforge.rendering.chart_renderers.placeholder import PlaceholderChartRenderer

        static_types = [
            "waterfall", "heatmap", "sankey", "gantt", "football_field",
            "sensitivity_table", "funnel", "treemap", "tornado", "sunburst",
        ]
        for chart_type in static_types:
            renderer = CHART_RENDERERS[chart_type]
            assert not isinstance(renderer, PlaceholderChartRenderer), (
                f"{chart_type} is still a PlaceholderChartRenderer"
            )
