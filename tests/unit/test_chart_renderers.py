"""Tests for chart renderers — native editable PPTX charts via python-pptx."""

from __future__ import annotations

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from deckforge.ir.charts.types import (
    BarChartData,
    StackedBarChartData,
    GroupedBarChartData,
    HorizontalBarChartData,
    LineChartData,
    MultiLineChartData,
    AreaChartData,
    StackedAreaChartData,
    PieChartData,
    DonutChartData,
    ChartDataSeries,
    WaterfallChartData,
)
from deckforge.ir.elements.base import Position
from deckforge.themes.types import (
    ResolvedTheme,
    ThemeColors,
    ThemeTypography,
    ThemeSpacing,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture
def pptx_slide():
    """Create a blank PPTX presentation slide for chart testing."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    return slide


@pytest.fixture
def theme():
    """Minimal resolved theme for chart formatting."""
    return ResolvedTheme(
        name="test-theme",
        description="Test theme",
        colors=ThemeColors(
            primary="#2E86AB",
            secondary="#A23B72",
            accent="#F18F01",
            background="#1A1A2E",
            surface="#16213E",
            text_primary="#EAEAEA",
            text_secondary="#B0B0B0",
            text_muted="#808080",
            positive="#4CAF50",
            negative="#F44336",
            warning="#FF9800",
        ),
        typography=ThemeTypography(
            heading_family="Arial",
            body_family="Calibri",
            mono_family="Consolas",
        ),
        spacing=ThemeSpacing(
            margin_top=0.5,
            margin_bottom=0.5,
            margin_left=0.75,
            margin_right=0.75,
            gutter=0.3,
            element_gap=0.2,
            section_gap=0.5,
        ),
        chart_colors=["#2E86AB", "#A23B72", "#F18F01", "#4CAF50", "#FF9800", "#9C27B0"],
    )


@pytest.fixture
def position():
    """Standard chart position on a slide."""
    return Position(x=1.0, y=1.5, width=11.0, height=5.0)


# ── Registry Tests ────────────────────────────────────────────────────────────


class TestChartRegistry:
    def test_chart_renderers_dict_exists(self):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        assert isinstance(CHART_RENDERERS, dict)

    def test_registry_contains_all_native_types(self):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        native_types = [
            "bar", "stacked_bar", "grouped_bar", "horizontal_bar",
            "line", "multi_line", "area", "stacked_area",
            "pie", "donut",
        ]
        for chart_type in native_types:
            assert chart_type in CHART_RENDERERS, f"Missing: {chart_type}"

    def test_registry_contains_placeholder_types(self):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        placeholder_types = [
            "waterfall", "funnel", "treemap", "tornado",
            "football_field", "sensitivity_table",
            "heatmap", "sankey", "gantt", "sunburst",
        ]
        for chart_type in placeholder_types:
            assert chart_type in CHART_RENDERERS, f"Missing placeholder: {chart_type}"

    def test_render_chart_function_exists(self):
        from deckforge.rendering.chart_renderers import render_chart

        assert callable(render_chart)

    def test_render_chart_dispatches_bar(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import render_chart

        chart_data = BarChartData(
            categories=["Q1", "Q2", "Q3"],
            series=[ChartDataSeries(name="Revenue", values=[100, 120, 140])],
            title="Revenue by Quarter",
        )
        render_chart(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) > 0


# ── Bar Chart Tests ───────────────────────────────────────────────────────────


class TestBarChartRenderer:
    def test_bar_chart_creates_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BarChartData(
            categories=["Q1", "Q2", "Q3", "Q4"],
            series=[ChartDataSeries(name="Revenue", values=[100, 120, 140, 160])],
            title="Revenue by Quarter",
        )
        renderer = CHART_RENDERERS["bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1

    def test_bar_chart_has_chart_object(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BarChartData(
            categories=["Q1", "Q2"],
            series=[ChartDataSeries(name="Rev", values=[100, 200])],
        )
        renderer = CHART_RENDERERS["bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_chart

    def test_stacked_bar_creates_chart(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = StackedBarChartData(
            categories=["A", "B"],
            series=[
                ChartDataSeries(name="S1", values=[10, 20]),
                ChartDataSeries(name="S2", values=[30, 40]),
            ],
        )
        renderer = CHART_RENDERERS["stacked_bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_chart

    def test_grouped_bar_creates_chart(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = GroupedBarChartData(
            categories=["A", "B"],
            series=[ChartDataSeries(name="S1", values=[10, 20])],
        )
        renderer = CHART_RENDERERS["grouped_bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_horizontal_bar_creates_chart(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = HorizontalBarChartData(
            categories=["X", "Y"],
            series=[ChartDataSeries(name="S1", values=[50, 60])],
        )
        renderer = CHART_RENDERERS["horizontal_bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_bar_chart_with_title(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BarChartData(
            categories=["Q1"],
            series=[ChartDataSeries(name="Rev", values=[100])],
            title="My Chart Title",
        )
        renderer = CHART_RENDERERS["bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "My Chart Title"

    def test_bar_chart_without_title(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BarChartData(
            categories=["Q1"],
            series=[ChartDataSeries(name="Rev", values=[100])],
        )
        renderer = CHART_RENDERERS["bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert not chart.has_title


# ── Line Chart Tests ──────────────────────────────────────────────────────────


class TestLineChartRenderer:
    def test_line_chart_creates_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = LineChartData(
            categories=["Jan", "Feb", "Mar"],
            series=[ChartDataSeries(name="Sales", values=[50, 60, 70])],
            title="Monthly Sales",
        )
        renderer = CHART_RENDERERS["line"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_multi_line_creates_chart(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = MultiLineChartData(
            categories=["Jan", "Feb"],
            series=[
                ChartDataSeries(name="S1", values=[10, 20]),
                ChartDataSeries(name="S2", values=[30, 40]),
            ],
        )
        renderer = CHART_RENDERERS["multi_line"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_line_chart_with_title(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = LineChartData(
            categories=["Q1"],
            series=[ChartDataSeries(name="Rev", values=[100])],
            title="Line Title",
        )
        renderer = CHART_RENDERERS["line"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert chart.has_title


# ── Area Chart Tests ──────────────────────────────────────────────────────────


class TestAreaChartRenderer:
    def test_area_chart_creates_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = AreaChartData(
            categories=["Jan", "Feb", "Mar"],
            series=[ChartDataSeries(name="S1", values=[10, 20, 30])],
        )
        renderer = CHART_RENDERERS["area"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_stacked_area_creates_chart(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = StackedAreaChartData(
            categories=["Jan", "Feb"],
            series=[
                ChartDataSeries(name="S1", values=[10, 20]),
                ChartDataSeries(name="S2", values=[30, 40]),
            ],
        )
        renderer = CHART_RENDERERS["stacked_area"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart


# ── Pie Chart Tests ───────────────────────────────────────────────────────────


class TestPieChartRenderer:
    def test_pie_chart_creates_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = PieChartData(
            labels=["A", "B", "C"],
            values=[50, 30, 20],
            title="Market Share",
        )
        renderer = CHART_RENDERERS["pie"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_pie_chart_has_title(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = PieChartData(
            labels=["A", "B"],
            values=[60, 40],
            title="Pie Title",
        )
        renderer = CHART_RENDERERS["pie"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "Pie Title"


# ── Donut Chart Tests ─────────────────────────────────────────────────────────


class TestDonutChartRenderer:
    def test_donut_chart_creates_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = DonutChartData(
            labels=["X", "Y", "Z"],
            values=[40, 35, 25],
        )
        renderer = CHART_RENDERERS["donut"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_donut_chart_is_doughnut_type(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS
        from pptx.enum.chart import XL_CHART_TYPE

        chart_data = DonutChartData(
            labels=["A", "B"],
            values=[60, 40],
        )
        renderer = CHART_RENDERERS["donut"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        # Verify it is a doughnut chart (not a pie)
        plot = chart.plots[0]
        assert plot.chart_type == XL_CHART_TYPE.DOUGHNUT


# ── Placeholder Tests ─────────────────────────────────────────────────────────


class TestPlaceholderChartRenderer:
    def test_placeholder_creates_shape(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = WaterfallChartData(
            categories=["Start", "Change", "End"],
            values=[100, -20, 80],
        )
        renderer = CHART_RENDERERS["waterfall"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) > 0

    def test_placeholder_has_text_label(self, pptx_slide, theme, position):
        from deckforge.rendering.chart_renderers import CHART_RENDERERS

        chart_data = WaterfallChartData(
            categories=["Start", "End"],
            values=[100, 80],
        )
        renderer = CHART_RENDERERS["waterfall"]
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_text_frame
        assert "waterfall" in shape.text_frame.text.lower()
